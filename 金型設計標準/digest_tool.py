#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""「プレス金型 設計・組立の基礎」 Excel編集 → PDF / PowerPoint

  ┌─ init ──→ Excel を作る（初回のみ。以後Excelが正）
  │
  Excelを編集（文言・数値・行の追加削除）
  │
  ├─ pdf ───→ A4縦 PDF を作る
  └─ pptx ──→ PowerPoint を作る（1トピック1スライド）

使い方:
    python digest_tool.py init  内容.xlsx          # 初回：Excelを作る
    python digest_tool.py pdf   内容.xlsx 出力.pdf
    python digest_tool.py pptx  内容.xlsx 出力.pptx

Excelの構成
  00_設定    … タイトル・出典など
  01_本文    … 第1部。種別(kind)で見た目が決まる
  02_トラブル … 第2部の相関図。グループ／図名／症状／確認／対策
  03_要点    … 最後のまとめ

種別(kind)の意味
  part      大見出し（部の扉）      h1  章見出し（紺帯）
  h2        小見出し                text  本文
  bullet    箇条書き                box  囲み（黄）      box_green 囲み（緑）
  table     表（本文に「||」で列、改行で行。列幅はoptionにmm指定）
  flow      横並びの流れ図          pagebreak 改ページ

依存: pip install openpyxl reportlab python-pptx
"""
import os
import sys

# ─────────────────────────── 共通 ───────────────────────────
NAVY_H = '1f3864'
TEAL_H = '1c6b63'
BLUE_H = '2e5c9a'

SHEETS = ('00_設定', '01_本文', '02_トラブル', '03_要点')


def _load(xlsx):
    """Excelから内容を読む。戻り値: (meta, body, diagrams, points)"""
    import openpyxl
    wb = openpyxl.load_workbook(xlsx)

    meta = {}
    ws = wb['00_設定']
    for r in range(2, ws.max_row + 1):
        k, v = ws.cell(r, 1).value, ws.cell(r, 2).value
        if k:
            meta[str(k).strip()] = str(v or '')

    body = []
    ws = wb['01_本文']
    for r in range(2, ws.max_row + 1):
        kind = ws.cell(r, 1).value
        if not kind:
            continue
        body.append((str(kind).strip(), str(ws.cell(r, 2).value or ''),
                     str(ws.cell(r, 3).value or '')))

    diagrams, cur = [], None
    ws = wb['02_トラブル']
    for r in range(2, ws.max_row + 1):
        grp = ws.cell(r, 1).value
        name = ws.cell(r, 2).value
        sym = ws.cell(r, 3).value
        if not (name or sym):
            continue
        if name and (cur is None or cur[1] != str(name).strip()):
            cur = (str(grp or (cur[0] if cur else '')).strip(), str(name).strip(), [])
            diagrams.append(cur)
        if sym is not None and cur is not None:
            cur[2].append((str(sym or ''), str(ws.cell(r, 4).value or ''),
                           str(ws.cell(r, 5).value or '')))

    points = []
    ws = wb['03_要点']
    for r in range(2, ws.max_row + 1):
        v = ws.cell(r, 2).value
        if v:
            points.append(str(v).strip())

    return meta, body, diagrams, points


# ─────────────────────────── init ───────────────────────────
def cmd_init(out):
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    import digest_data as D

    WRAP = Alignment(wrap_text=True, vertical='top')
    HFILL = PatternFill('solid', fgColor='FF' + BLUE_H)
    HFONT = Font(name='Meiryo UI', size=9, bold=True, color='FFFFFFFF')
    BFONT = Font(name='Meiryo UI', size=9)

    def head(ws, cols, widths):
        for i, (c, w) in enumerate(zip(cols, widths), 1):
            cell = ws.cell(1, i, c)
            cell.fill, cell.font = HFILL, HFONT
            cell.alignment = Alignment(wrap_text=True, vertical='center', horizontal='center')
            ws.column_dimensions[get_column_letter(i)].width = w
        ws.freeze_panes = 'A2'

    wb = Workbook()

    ws = wb.active
    ws.title = '00_設定'
    head(ws, ['項目', '内容'], [18, 110])
    for r, (k, v) in enumerate([('タイトル', D.TITLE), ('サブタイトル', D.SUBTITLE),
                                ('前書き', D.INTRO), ('出典', D.SOURCE)], 2):
        ws.cell(r, 1, k).font = BFONT
        c = ws.cell(r, 2, v)
        c.font, c.alignment = BFONT, WRAP
        ws.row_dimensions[r].height = 46

    ws = wb.create_sheet('01_本文')
    head(ws, ['種別', 'option', '内容'], [12, 14, 120])
    for r, (kind, opt, txt) in enumerate(D.BODY, 2):
        ws.cell(r, 1, kind).font = BFONT
        ws.cell(r, 2, opt).font = BFONT
        c = ws.cell(r, 3, txt)
        c.font, c.alignment = BFONT, WRAP
        ws.row_dimensions[r].height = max(18, min(120, len(txt) * 0.42))

    ws = wb.create_sheet('02_トラブル')
    head(ws, ['グループ', '図番号・名前', '症状・条件', '確認すること', '対策'], [22, 26, 34, 34, 62])
    r = 2
    for grp, name, rows in D.DIAGRAMS:
        for i, (a, b, c) in enumerate(rows):
            ws.cell(r, 1, grp if i == 0 else '').font = BFONT
            ws.cell(r, 2, name if i == 0 else '').font = BFONT
            for j, v in enumerate((a, b, c), 3):
                cc = ws.cell(r, j, v)
                cc.font, cc.alignment = BFONT, WRAP
            ws.row_dimensions[r].height = max(16, min(60, len(c) * 0.5))
            r += 1

    ws = wb.create_sheet('03_要点')
    head(ws, ['No', '内容'], [6, 110])
    for r, t in enumerate(D.POINTS, 2):
        ws.cell(r, 1, r - 1).font = BFONT
        c = ws.cell(r, 2, t)
        c.font, c.alignment = BFONT, WRAP

    wb.save(out)
    print('作成しました:', out)
    print('  以後はこのExcelを編集して pdf / pptx を作り直してください')


# ─────────────────────────── PDF ───────────────────────────
def cmd_pdf(xlsx, out):
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.platypus import (HRFlowable, KeepTogether, PageBreak, Paragraph,
                                    SimpleDocTemplate, Spacer, Table, TableStyle)

    meta, body, diagrams, points = _load(xlsx)
    pdfmetrics.registerFont(UnicodeCIDFont('HeiseiKakuGo-W5'))
    F = 'HeiseiKakuGo-W5'
    NAVY = colors.HexColor('#' + NAVY_H)
    TEAL = colors.HexColor('#' + TEAL_H)
    BLUE = colors.HexColor('#' + BLUE_H)
    MGRAY = colors.HexColor('#c8cfda')
    LGRAY = colors.HexColor('#f2f4f8')
    W = 188 * mm

    S = {
        'title': ParagraphStyle('t', fontName=F, fontSize=17, textColor=NAVY, alignment=TA_CENTER, leading=22),
        'sub':   ParagraphStyle('s', fontName=F, fontSize=8.5, textColor=colors.HexColor('#555'),
                                alignment=TA_CENTER, leading=12),
        'part':  ParagraphStyle('pt', fontName=F, fontSize=14, textColor=colors.white, backColor=TEAL,
                                leftIndent=6, spaceBefore=4, spaceAfter=6, leading=22),
        'h1':    ParagraphStyle('h1', fontName=F, fontSize=12, textColor=colors.white, backColor=NAVY,
                                leftIndent=5, spaceBefore=8, spaceAfter=4, leading=18),
        'h2':    ParagraphStyle('h2', fontName=F, fontSize=10.5, textColor=NAVY,
                                spaceBefore=6, spaceAfter=2, leading=15),
        'b':     ParagraphStyle('b', fontName=F, fontSize=9, leading=13.8),
        'sm':    ParagraphStyle('sm', fontName=F, fontSize=8.3, leading=12.5),
        'hw':    ParagraphStyle('hw', fontName=F, fontSize=8.8, textColor=colors.white, leading=13),
        'hwc':   ParagraphStyle('hwc', fontName=F, fontSize=8.8, textColor=colors.white,
                                leading=13, alignment=TA_CENTER),
        'note':  ParagraphStyle('n', fontName=F, fontSize=7.8, textColor=colors.HexColor('#666'), leading=11),
    }
    P = lambda t, s='b': Paragraph(t, S[s])

    def mk_table(rows, widths, head=True, bg=None, small=True):
        t = Table(rows, colWidths=widths, repeatRows=1 if head else 0)
        st = [('GRID', (0, 0), (-1, -1), 0.5, MGRAY), ('VALIGN', (0, 0), (-1, -1), 'TOP'),
              ('TOPPADDING', (0, 0), (-1, -1), 3.5), ('BOTTOMPADDING', (0, 0), (-1, -1), 3.5),
              ('LEFTPADDING', (0, 0), (-1, -1), 5)]
        if head:
            st.append(('BACKGROUND', (0, 0), (-1, 0), BLUE))
        if bg:
            st.append(('BACKGROUND', (0, 1), (-1, -1), bg))
        t.setStyle(TableStyle(st))
        return t

    def mk_box(text, green=False):
        bg = colors.HexColor('#eef6ee' if green else '#fdf6e3')
        bd = colors.HexColor('#9fbf9f' if green else '#e0c66b')
        t = Table([[P(text)]], colWidths=[W])
        t.setStyle(TableStyle([('BACKGROUND', (0, 0), (-1, -1), bg), ('BOX', (0, 0), (-1, -1), 0.8, bd),
                               ('TOPPADDING', (0, 0), (-1, -1), 6), ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                               ('LEFTPADDING', (0, 0), (-1, -1), 8), ('RIGHTPADDING', (0, 0), (-1, -1), 8)]))
        return t

    doc = SimpleDocTemplate(out, pagesize=A4, leftMargin=11 * mm, rightMargin=11 * mm,
                            topMargin=10 * mm, bottomMargin=12 * mm,
                            title=meta.get('タイトル', ''), author='橋本工業株式会社')
    st = [P(meta.get('タイトル', ''), 'title'), P(meta.get('サブタイトル', ''), 'sub'),
          HRFlowable(width='100%', thickness=1.5, color=NAVY, spaceAfter=6),
          mk_box(meta.get('前書き', '')), Spacer(1, 3 * mm)]

    for kind, opt, txt in body:
        if kind == 'part':
            st.append(P(txt, 'part'))
        elif kind == 'h1':
            st.append(P('■ ' + txt, 'h1'))
        elif kind == 'h2':
            st.append(P('◆ ' + txt, 'h2'))
        elif kind == 'text':
            st.append(P(txt))
        elif kind == 'bullet':
            st.append(P('・' + txt))
        elif kind == 'box':
            st.append(mk_box(txt))
        elif kind == 'box_green':
            st.append(mk_box(txt, green=True))
        elif kind == 'pagebreak':
            st.append(PageBreak())
        elif kind == 'flow':
            cells = [c.replace('\n', '<br/>') for c in txt.split('||')]
            n = (len(cells) + 1) // 2
            top, bottom = cells[:n], cells[n:]
            bottom += [''] * (n - len(bottom))
            rows = [[P(c, 'hwc') for c in top], [P(c, 'sm') for c in bottom]]
            st.append(mk_table(rows, [W / n] * n))
        elif kind == 'table':
            widths = [float(x) * mm for x in opt.split(',')] if opt else None
            lines = [l for l in txt.split('\n') if l.strip()]
            rows = []
            for i, line in enumerate(lines):
                cs = line.split('||')
                rows.append([P(c, 'hw' if i == 0 else 'sm') for c in cs])
            if not widths:
                widths = [W / len(rows[0])] * len(rows[0])
            st.append(mk_table(rows, widths))

    # 第2部
    prev_grp = None
    for grp, name, rows in diagrams:
        blk = []
        if grp != prev_grp:
            if prev_grp is not None:
                st.append(PageBreak())
            st.append(P('■ ' + grp, 'h1'))
            prev_grp = grp
        data = [[P('<b>症状・条件</b>', 'hw'), P('<b>確認すること</b>', 'hw'), P('<b>対策</b>', 'hw')]]
        for a, b, c in rows:
            data.append([P(a, 'sm'), P(b, 'sm'), P(c, 'sm')])
        blk.append(P('◆ ' + name, 'h2'))
        blk.append(mk_table(data, [46 * mm, 62 * mm, 80 * mm]))
        st.append(KeepTogether(blk))

    st.append(P('■ 覚えておきたい要点', 'h1'))
    st.append(mk_table([[P('<b>%d</b>' % i, 'sm'), P(t, 'sm')] for i, t in enumerate(points, 1)],
                       [12 * mm, 176 * mm], head=False, bg=LGRAY))
    st.append(Spacer(1, 3 * mm))
    st.append(P(meta.get('出典', ''), 'note'))

    doc.build(st)
    print('作成しました:', out)


# ─────────────────────────── PPTX ───────────────────────────
def cmd_pptx(xlsx, out):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    meta, body, diagrams, points = _load(xlsx)
    NAVY = RGBColor(0x1F, 0x38, 0x64)
    TEAL = RGBColor(0x1C, 0x6B, 0x63)
    BLUE = RGBColor(0x2E, 0x5C, 0x9A)
    GRAY = RGBColor(0x44, 0x44, 0x44)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    FONT = 'Meiryo UI'

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    SW, SH = 13.333, 7.5
    blank = prs.slide_layouts[6]

    def add_slide():
        return prs.slides.add_slide(blank)

    def txbox(sl, x, y, w, h, text, size=18, bold=False, color=GRAY, align=None):
        tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]
        p.text = text
        if align is not None:
            p.alignment = align
        for r in p.runs:
            r.font.name, r.font.size, r.font.bold, r.font.color.rgb = FONT, Pt(size), bold, color
        return tb

    def band(sl, text, color=NAVY, y=0.35, size=26):
        sh = sl.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(SW - 0.8), Inches(0.85))
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        sh.line.fill.background()
        tf = sh.text_frame
        tf.margin_left = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = text
        for r in p.runs:
            r.font.name, r.font.size, r.font.bold, r.font.color.rgb = FONT, Pt(size), True, WHITE

    def table(sl, rows, widths, top, head=True, fsize=12):
        nr = len(rows)
        nc = max(len(r) for r in rows)
        # 列数が足りない行を空セルで埋める（不揃いだと生成時に落ちるため）
        rows = [list(r) + [''] * (nc - len(r)) for r in rows]
        widths = (list(widths) + [1] * nc)[:nc]
        tb = sl.shapes.add_table(nr, nc, Inches(0.5), Inches(top),
                                 Inches(SW - 1.0), Inches(min(0.36 * nr, SH - top - 0.4))).table
        total = sum(widths)
        for i, w in enumerate(widths):
            tb.columns[i].width = Inches((SW - 1.0) * w / total)
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                cell = tb.cell(ri, ci)
                cell.text = val
                cell.margin_left = cell.margin_right = Inches(0.06)
                cell.margin_top = cell.margin_bottom = Inches(0.02)
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.name = FONT
                        r.font.size = Pt(fsize)
                        r.font.bold = (ri == 0 and head)
                        r.font.color.rgb = WHITE if (ri == 0 and head) else GRAY
                if ri == 0 and head:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = BLUE
        return tb

    # 表紙
    sl = add_slide()
    txbox(sl, 1.0, 2.4, SW - 2.0, 1.2, meta.get('タイトル', ''), 40, True, NAVY)
    txbox(sl, 1.0, 3.7, SW - 2.0, 0.6, meta.get('サブタイトル', ''), 16, False, GRAY)
    txbox(sl, 1.0, 4.6, SW - 2.0, 1.4, meta.get('前書き', ''), 12, False, GRAY)

    # 第1部：h1ごとに1枚（内容が多ければ続きスライド）
    cur, buf = None, []

    def flush():
        """溜めた内容を1枚に出す"""
        if cur is None:
            return
        sl = add_slide()
        band(sl, cur)
        y = 1.45
        for kind, opt, txt in buf:
            if y > SH - 0.7:
                break
            if kind == 'h2':
                txbox(sl, 0.55, y, SW - 1.1, 0.4, '◆ ' + txt, 17, True, NAVY)
                y += 0.5
            elif kind == 'bullet':
                n = txbox(sl, 0.75, y, SW - 1.4, 0.4, '・' + txt, 12)
                y += 0.34 + 0.16 * (len(txt) // 60)
            elif kind == 'text':
                txbox(sl, 0.55, y, SW - 1.1, 0.4, txt, 12)
                y += 0.34 + 0.16 * (len(txt) // 60)
            elif kind in ('box', 'box_green'):
                sh = sl.shapes.add_shape(1, Inches(0.55), Inches(y), Inches(SW - 1.1), Inches(0.62))
                sh.fill.solid()
                sh.fill.fore_color.rgb = RGBColor(0xEE, 0xF6, 0xEE) if kind == 'box_green' \
                    else RGBColor(0xFD, 0xF6, 0xE3)
                sh.line.color.rgb = RGBColor(0x9F, 0xBF, 0x9F) if kind == 'box_green' \
                    else RGBColor(0xE0, 0xC6, 0x6B)
                tf = sh.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.15)
                p = tf.paragraphs[0]
                p.text = txt
                for r in p.runs:
                    r.font.name, r.font.size, r.font.color.rgb = FONT, Pt(12), GRAY
                y += 0.78
            elif kind == 'table':
                lines = [l for l in txt.split('\n') if l.strip()]
                rows = [l.split('||') for l in lines]
                widths = [float(x) for x in opt.split(',')] if opt else [1] * len(rows[0])
                table(sl, rows, widths, y, fsize=11)
                y += 0.36 * len(rows) + 0.3
            elif kind == 'flow':
                # 1行目=見出し、2行目=説明。改行が入っていれば行区切りとして扱う
                lines = [l for l in txt.split('\n\n') if l.strip()] or [txt]
                if len(lines) == 1:
                    cells = txt.replace('\n', '@@').split('||')
                    n = (len(cells) + 1) // 2
                    rows = [[c.replace('@@', ' ') for c in cells[:n]],
                            [c.replace('@@', ' ') for c in cells[n:]]]
                else:
                    rows = [[c.replace('\n', ' ') for c in l.split('||')] for l in lines]
                table(sl, rows, [1] * max(len(r) for r in rows), y, fsize=11)
                y += 1.2
        buf.clear()

    for kind, opt, txt in body:
        if kind == 'part':
            flush()
            cur = None
            sl = add_slide()
            band(sl, txt, TEAL, y=3.1, size=30)
        elif kind == 'h1':
            flush()
            cur = txt
        elif kind == 'pagebreak':
            continue
        else:
            buf.append((kind, opt, txt))
    flush()

    # 第2部：1相関図＝1スライド
    prev = None
    for grp, name, rows in diagrams:
        if grp != prev:
            sl = add_slide()
            band(sl, grp, TEAL, y=3.1, size=30)
            prev = grp
        sl = add_slide()
        band(sl, name)
        data = [['症状・条件', '確認すること', '対策']] + [list(r) for r in rows]
        table(sl, data, [24, 30, 46], 1.45, fsize=11 if len(rows) <= 7 else 10)

    # 要点
    sl = add_slide()
    band(sl, '覚えておきたい要点')
    table(sl, [[str(i), t] for i, t in enumerate(points, 1)], [4, 96], 1.5, head=False, fsize=14)

    sl = add_slide()
    txbox(sl, 1.0, 3.0, SW - 2.0, 1.5, meta.get('出典', ''), 12, False, GRAY)

    prs.save(out)
    print('作成しました:', out, '（%dスライド）' % len(prs.slides._sldIdLst))


# ─────────────────────────── main ───────────────────────────
if __name__ == '__main__':
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)
    cmd = sys.argv[1]
    if cmd == 'init':
        cmd_init(sys.argv[2])
    elif cmd == 'pdf':
        cmd_pdf(sys.argv[2], sys.argv[3] if len(sys.argv) > 3 else 'digest.pdf')
    elif cmd == 'pptx':
        cmd_pptx(sys.argv[2], sys.argv[3] if len(sys.argv) > 3 else 'digest.pptx')
    else:
        print(__doc__)
        sys.exit(1)
