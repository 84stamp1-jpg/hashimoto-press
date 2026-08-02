#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""HMDS金型設計標準の3編を、編集可能なPowerPoint（文書レイアウト）にする。

build_hms_docs.py と同じ内容（Excel＋クリアランス基準書＋HMS図）を読み込み、
PDFと同じ体裁を A4縦のスライドに流し込む。1スライド＝1ページ相当。
図は画像として貼るので、PowerPoint上で外周トリミング・被せ図形の微調整ができる。

使い方:
    python build_hms_pptx.py            # 3編すべて pptx/ へ
    python build_hms_pptx.py 共通
依存: pip install python-pptx openpyxl pillow
"""
import math
import os
import re
import sys

from PIL import Image
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import build_hms_docs as B   # 内容ローダーを再利用

FONT = 'Meiryo'
NAVY = '1F3864'
TEAL = '1C6B63'
GRAY = '555555'
BLACK = '111111'
LINE = 'B7C2D0'
LIGHTBG = 'F5F8FB'
HEADBG = 'EAF0F6'

# A4縦とレイアウト（cm）
PW, PH = 21.0, 29.7
ML, MR, MT, MB = 1.8, 1.8, 1.6, 1.7
CW = PW - ML - MR                 # 本文幅
BOTTOM = PH - MB                   # 下限y
CPL = 43                           # 本文1行あたりの全角文字数の目安
LH = 0.52                          # 本文の行送り(cm)


def _set_font(run, size, color, bold=False, name=FONT):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.name = name
    f.color.rgb = RGBColor.from_string(color)
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)


def _lines(text):
    n = 0
    for seg in str(text).split('\n'):
        n += max(1, math.ceil(len(seg) / CPL))
    return n


class Doc:
    def __init__(self, prs, title_ja):
        self.prs = prs
        self.title_ja = title_ja
        self.blank = prs.slide_layouts[6]
        self.slide = None
        self.y = MT

    # ---- スライド管理 -------------------------------------------------
    def new_slide(self, footer=True):
        self.slide = self.prs.slides.add_slide(self.blank)
        self.y = MT
        if footer:
            self._footer()
        return self.slide

    def _footer(self):
        tb = self.slide.shapes.add_textbox(Cm(ML), Cm(PH - 1.25), Cm(CW), Cm(0.7))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = 'HMDS 金型設計標準【%s】　橋本工業株式会社' % self.title_ja
        _set_font(r, 8, GRAY)
        # ページ番号
        tb2 = self.slide.shapes.add_textbox(Cm(PW - MR - 2), Cm(PH - 1.25), Cm(2), Cm(0.7))
        tf2 = tb2.text_frame
        tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        r2.text = str(len(self.prs.slides._sldIdLst))
        _set_font(r2, 8, GRAY)
        # 罫線
        ln = self.slide.shapes.add_connector(2, Cm(ML), Cm(PH - 1.35),
                                             Cm(PW - MR), Cm(PH - 1.35))
        ln.line.color.rgb = RGBColor.from_string(LINE)
        ln.line.width = Pt(0.5)

    def ensure(self, h):
        if self.slide is None or self.y + h > BOTTOM:
            self.new_slide()

    # ---- 要素 ---------------------------------------------------------
    def bar(self, text, page_break=False):
        if page_break or self.slide is None or self.y + 1.1 > BOTTOM:
            self.new_slide()
        h = 0.9
        sp = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(ML), Cm(self.y),
                                         Cm(CW), Cm(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(NAVY)
        sp.line.fill.background()
        sp.shadow.inherit = False
        tf = sp.text_frame
        tf.margin_left = Cm(0.25)
        tf.margin_top = Cm(0.05)
        tf.margin_bottom = Cm(0.05)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        r = tf.paragraphs[0].add_run()
        r.text = text
        _set_font(r, 14, 'FFFFFF', bold=True)
        self.y += h + 0.35

    def _text(self, text, size, color, bold=False, indent=0.0, gap=0.18, teal=False):
        h = _lines(text) * LH + 0.12
        self.ensure(h)
        tb = self.slide.shapes.add_textbox(Cm(ML + indent), Cm(self.y),
                                           Cm(CW - indent), Cm(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        first = True
        for seg in str(text).split('\n'):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.line_spacing = 1.15
            r = p.add_run()
            r.text = seg
            _set_font(r, size, color, bold=bold)
        self.y += h + gap
        return tb

    def heading(self, no, item):
        self._text('%s　%s' % (no, item), 11, NAVY, bold=True, gap=0.06)

    def body(self, text):
        self._text(text, 10.5, BLACK, gap=0.2)

    def ref(self, text):
        self._text('▶ ' + text, 9.5, TEAL, gap=0.12)

    def note_intro(self, text):
        self._text(text, 10, BLACK, gap=0.15)

    # 小見出し類（早見表用）
    def sub(self, text):
        self._text(text, 9.5, GRAY, gap=0.1)

    def head2(self, text):
        self._text('■ ' + text, 10.5, TEAL, bold=True, gap=0.1)

    def note(self, text):
        self._text('※ ' + text, 8.5, GRAY, gap=0.1)

    def title2(self, text):
        self._text('◆ ' + text, 10.5, NAVY, bold=True, gap=0.08)

    def figure(self, path):
        iw, ih = Image.open(path).size
        natural_w = iw / 47.0                  # ≒120dpi
        w = min(CW, natural_w)
        h = w * ih / iw
        max_h = BOTTOM - MT - 1.0
        if h > max_h:
            h = max_h
            w = h * iw / ih
        # 収まらなければ次スライドへ
        if self.y + h + 0.2 > BOTTOM:
            self.new_slide()
        x = ML + (CW - w) / 2
        self.slide.shapes.add_picture(path, Cm(x), Cm(self.y), Cm(w), Cm(h))
        self.y += h + 0.3

    def table(self, rows, first_is_header=True):
        nr = len(rows)
        nc = max(len(r) for r in rows)
        rh = 0.72
        th = nr * rh
        if self.y + th > BOTTOM:
            self.new_slide()
        gtbl = self.slide.shapes.add_table(nr, nc, Cm(ML), Cm(self.y),
                                           Cm(CW), Cm(th)).table
        # 列幅：1列目広め
        first_w = CW * 0.26 if nc > 1 else CW
        rest = (CW - first_w) / (nc - 1) if nc > 1 else 0
        gtbl.columns[0].width = Cm(first_w)
        for c in range(1, nc):
            gtbl.columns[c].width = Cm(rest)
        for ri, row in enumerate(rows):
            gtbl.rows[ri].height = Cm(rh)
            for ci in range(nc):
                cell = gtbl.cell(ri, ci)
                cell.margin_left = Cm(0.12)
                cell.margin_right = Cm(0.08)
                cell.margin_top = Cm(0.02)
                cell.margin_bottom = Cm(0.02)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                val = row[ci] if ci < len(row) else ''
                is_h = first_is_header and ri == 0
                cell.fill.solid()
                if is_h:
                    cell.fill.fore_color.rgb = RGBColor.from_string(TEAL)
                else:
                    cell.fill.fore_color.rgb = RGBColor.from_string(
                        LIGHTBG if ri % 2 else 'FFFFFF')
                para = cell.text_frame.paragraphs[0]
                r = para.add_run()
                r.text = str(val)
                _set_font(r, 8.5, 'FFFFFF' if is_h else BLACK, bold=is_h)
        self.y += th + 0.3


# ================================================================ 組み立て
def _clean_rule(rule):
    rule = re.sub(r'※?\s*(図|表|写真)?は?他社資料\s*S\d+[〜～\-]*S?\d*[^。]*?(参照|あり)。?',
                  '', rule)
    rule = re.sub(r'（\s*S\d+[〜～\-]*S?\d*[^）]*）', '', rule)
    rule = re.sub(r'※?\s*S\d+\s*に写真あり。?', '', rule)
    return rule.strip()


def render_blocks(doc, blocks):
    for kind, payload in blocks:
        if kind == 'title':
            doc.title2(payload)
        elif kind == 'sub':
            doc.sub(payload)
        elif kind == 'head':
            doc.head2(payload)
        elif kind == 'note':
            doc.note(payload)
        elif kind == 'para':
            doc.body(payload)
        elif kind == 'table':
            doc.table(payload)


def build_cover(doc, title_ja, intro):
    doc.new_slide(footer=False)
    s = doc.slide
    # 上部の帯
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Cm(PW), Cm(2.5))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor.from_string(NAVY)
    band.line.fill.background()
    band.shadow.inherit = False
    # タイトル
    tb = s.shapes.add_textbox(Cm(ML), Cm(9), Cm(CW), Cm(4))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, sz) in enumerate([('HMDS 金型設計標準', 30), ('【%s】' % title_ja, 30)]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = txt
        _set_font(r, sz, NAVY, bold=True)
    tb2 = s.shapes.add_textbox(Cm(ML), Cm(14), Cm(CW), Cm(2))
    tf2 = tb2.text_frame
    for i, (txt, sz, col) in enumerate([('橋本工業株式会社　技術部', 13, BLACK),
                                        ('HASHIMOTO KOGYO — Mold Design Standard（HMDS）', 11, GRAY)]):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        r = p.add_run()
        r.text = txt
        _set_font(r, sz, col)
    # はじめに表
    rows = [(a, b) for a, b in intro if a or b]
    if rows:
        y = 17.5
        tbl = s.shapes.add_table(len(rows), 2, Cm(ML), Cm(y),
                                 Cm(CW), Cm(0.7 * len(rows))).table
        tbl.columns[0].width = Cm(4.5)
        tbl.columns[1].width = Cm(CW - 4.5)
        for ri, (a, b) in enumerate(rows):
            for ci, val in enumerate((a, b)):
                cell = tbl.cell(ri, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string('FFFFFF')
                cell.margin_top = cell.margin_bottom = Cm(0.03)
                p = cell.text_frame.paragraphs[0]
                p.line_spacing = 1.1
                r = p.add_run()
                r.text = val
                _set_font(r, 9.5, NAVY if ci == 0 else BLACK, bold=(ci == 0))


def build_edition(prs, key, clear_section):
    fname, sheet, ja, prefix = B.EDITIONS[key]
    path = os.path.join(B.BASE, fname)
    chapters = B.load_edition(path, sheet)
    intro = B.load_intro(path)
    doc = Doc(prs, ja)
    build_cover(doc, ja, intro)
    doc.new_slide()

    for cname, items in chapters:
        pb = '送り・ガイド' in cname
        doc.bar(cname, page_break=pb)
        for it in items:
            if it['kind'] == 'ref':
                doc.ref(it['text'])
                continue
            if it['dec'] in ('対象外', '不要'):
                continue
            body = _clean_rule(it['rule'])
            fig = os.path.join(B.FIGDIR, 'hms_%s%s.png' % (prefix, it['no']))
            has_fig = os.path.exists(fig)
            # 見出し＋本文＋図がなるべく同一スライドに載るよう、必要なら改ページ
            need = _lines('%s %s' % (it['no'], it['item'])) * LH + 0.2
            need += _lines(body) * LH + 0.3
            if has_fig:
                iw, ih = Image.open(fig).size
                w = min(CW, iw / 47.0)
                need += min(w * ih / iw, BOTTOM - MT - 1.0) + 0.3
            if doc.y + need > BOTTOM and need <= (BOTTOM - MT):
                doc.new_slide()
            doc.heading(it['no'], it['item'])
            if body:
                doc.body(body)
            if has_fig:
                doc.figure(fig)
        if clear_section and '設計値' in cname:
            bar_title, intro_text, sheets = clear_section
            doc.bar(bar_title)
            doc.note_intro(intro_text)
            for sh, blocks in sheets:
                render_blocks(doc, blocks)


def main():
    outdir = os.path.join(B.BASE, '出力')
    os.makedirs(outdir, exist_ok=True)
    keys = [sys.argv[1]] if len(sys.argv) > 1 else list(B.EDITIONS)
    cl = B.load_clearance()

    def sheets(names):
        return [(n, cl[n]) for n in names if n in cl]

    common_clear = ('5-補. クリアランス・曲げ 早見表',
                    'クリアランス（ピアス・抜き／バーリング）と曲げの早見表。'
                    '数値は標準値で、材料ロット・型構造・製品要求により調整する。',
                    sheets(['ピアス・抜き', 'バーリング', '曲げ']))
    junso_clear = ('補. 順送レイアウト早見表',
                   '順送型のレイアウト設計基準（早見表）。数値は標準値で、'
                   '製品・材料により調整する。',
                   sheets(['順送レイアウト']))
    CLEAR = {'共通': common_clear, '単発': None, '順送': junso_clear}

    for k in keys:
        prs = Presentation()
        prs.slide_width = Cm(PW)
        prs.slide_height = Cm(PH)
        build_edition(prs, k, CLEAR[k])
        ja = B.EDITIONS[k][2]
        out = os.path.join(outdir, 'HMDS金型設計標準_%s.pptx' % ja)
        prs.save(out)
        print('PPTX:', out)


if __name__ == '__main__':
    main()
